# PPTX (PowerPoint) 파일 파싱을 위한 라이브러리
from pptx import Presentation

# PPTX 파일 파싱 함수
def parse_pptx(file_path):
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                # 텍스트 상자 내용 출력
                print(f"PPTX Text Box: {text}")

# 파일 경로 지정
file_path = "example.pptx"  # 예제 PPTX 파일의 경로로 수정

if file_path.endswith(".pptx"):
    parse_pptx(file_path)
